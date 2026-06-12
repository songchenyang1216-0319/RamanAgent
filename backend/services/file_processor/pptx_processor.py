from __future__ import annotations

from pathlib import Path

from .base import BaseFileProcessor, ProcessedFile


class PptxFileProcessor(BaseFileProcessor):
    file_type = "pptx"
    supported_suffixes = {".pptx"}

    def process(self, path: Path, *, file_id: str | None = None, **_: object) -> ProcessedFile:
        try:
            from pptx import Presentation
        except ModuleNotFoundError:
            return self.failure(path, "当前环境缺少 python-pptx，暂时无法提取 PPTX 正文。")
        try:
            presentation = Presentation(str(path))
            slide_texts = []
            empty_slides = []
            notes_chars = 0
            for slide_index, slide in enumerate(presentation.slides, start=1):
                parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text.strip())
                try:
                    notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
                except Exception:
                    notes = ""
                notes = str(notes or "").strip()
                if notes:
                    notes_chars += len(notes)
                    parts.append(f"Notes:\n{notes}")
                if parts:
                    slide_texts.append(f"# Slide {slide_index}\n" + "\n".join(parts))
                else:
                    empty_slides.append(slide_index)
            text = "\n\n".join(slide_texts)
        except Exception as exc:
            return self.failure(path, f"PPTX 文件解析失败：{exc}")
        return ProcessedFile(
            success=True,
            file_type=self.file_type,
            filename=path.name,
            summary=f"已解析 PPTX，共 {len(presentation.slides)} 页，提取到 {len(slide_texts)} 页文本。",
            metadata={
                "slides": len(presentation.slides),
                "slides_with_text": len(slide_texts),
                "empty_slides": empty_slides[:50],
                "empty_slide_count": len(empty_slides),
                "notes_characters": notes_chars,
                "characters": len(text),
                "suffix": path.suffix.lower(),
            },
            preview=text[:3000],
            chunks=self.make_chunks(text, file_id=file_id, section="pptx"),
        )
